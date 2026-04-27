#!/usr/bin/env python3
"""
pptx-audit: Deterministic 14-check audit engine for PowerPoint presentations.

Reads a .pptx file with python-pptx + lxml, runs 14 structural checks,
emits a JSON report to stdout, and exits 1 if any CRITICAL issues are
found (0 otherwise).

Designed to be called by an agent before delivering any pptx build:
    python3 audit.py /path/to/deck.pptx
    python3 audit.py /path/to/deck.pptx --style STYLE-02
    python3 audit.py /path/to/deck.pptx --json-only      # suppress stderr summary

Exit codes:
    0 — no CRITICAL issues
    1 — CRITICAL issues present
    2 — script-level error (file not found, bad pptx, missing deps)

Output (stdout) is always single-document JSON:
{
  "path": "/abs/path/to/deck.pptx",
  "style": "STYLE-02" or null,
  "summary": {"critical": N, "warning": N, "info": N, "passed": bool},
  "critical": [{"check": int, "slide": int, "shape": str, "msg": str, ...}, ...],
  "warning": [...],
  "info": [...]
}

Coverage and simplifications (read this before relying on results):
  ✓ CHECK 1 BOUNDS                 — full
  ✓ CHECK 2 TEXT CLIPPING          — full (simulate_wrap)
  ✓ CHECK 3 WORD-WRAP QUALITY      — full
  ✓ CHECK 4 CONTAINER-TEXT SYNC    — full (center-point containment)
  ~ CHECK 5 BULLET ALIGNMENT       — minimal: dot-shape detection only;
                                      strict line-count layout check skipped to avoid FPs
  ✓ CHECK 6 OVERLAP                — full
  ~ CHECK 7 Z-ORDER                — minimal: opaque-fill-over-text by area
  ✓ CHECK 8 FONT COMPLIANCE        — full (size + name explicitness)
  ~ CHECK 9 SPACING CONSISTENCY    — minimal: cross-slide left-margin variance
  ~ CHECK 10 COLOR/FILL INTEGRITY  — minimal: orphan accent1 theme fills only
  ~ CHECK 11 STYLE COMPLIANCE      — only when --style is passed; checks fonts + palette
                                      against the style dict in style-pptx-mapping.md
  ✓ CHECK 12 IMAGE AR DISTORTION   — full
  ✓ CHECK 13 BROKEN GRADIENTS      — full (the "blue rectangle" bug)
  ~ CHECK 14 TEXT ZONE LUMINANCE   — full when shape.name carries the
                                      "text_zone=ZONE:SIZE:COLOR" annotation

A "~" check is conservative — it flags only high-confidence cases to
avoid the audit-credibility-killing false positives the original
audit-system.md warns about. The full spec lives in
references/audit-checks.md.
"""

import argparse
import io
import json
import os
import re
import sys
import tempfile
from collections import defaultdict
from pathlib import Path

# ----------------------------------------------------------------------
# Constants
# ----------------------------------------------------------------------

# Default 16:9 EMUs; we read actual dims from the deck.
EMU_PER_INCH = 914400
EMU_PER_PT = 12700

# CHECK 1 tolerance (~0.016")
BOUNDS_TOL = 15000

# CHECK 3 char-width factors (multiply by font_size_emu)
CHAR_WIDTH_FACTORS = {
    ('Georgia', False):       0.57, ('Georgia', True):       0.62,
    ('Rockwell', False):      0.60, ('Rockwell', True):      0.62,
    ('Calibri', False):       0.55, ('Calibri', True):       0.58,
    ('Arial', False):         0.55, ('Arial', True):         0.58,
    ('Helvetica', False):     0.55, ('Helvetica', True):     0.58,
    ('Times New Roman', False): 0.55, ('Times New Roman', True): 0.60,
    ('Consolas', False):      0.60, ('Consolas', True):      0.60,
    ('Arial Black', False):   0.63, ('Arial Black', True):   0.63,
    ('Comic Sans MS', False): 0.59, ('Comic Sans MS', True): 0.63,
}
DEFAULT_CHAR_FACTOR = 0.58
WORD_SAFETY = 1.12

# CHECK 2 line-height ratio
LINE_HEIGHT_RATIO = 1.35
TEXT_OVERFLOW_TOL = 1.10  # only flag when estimate > 110% of frame

# CHECK 8 minimum font size
MIN_FONT_PT = 14
MIN_FONT_PT_CAPTION = 10  # for caption-only, but we report below 14 as warning

# CHECK 12
AR_CRITICAL = 0.10
AR_WARNING = 0.05

# CHECK 14
LUM_CRITICAL_BRIGHT = 140  # white text needs zone <= this
LUM_CRITICAL_DARK = 115    # dark text needs zone >= this

# ----------------------------------------------------------------------
# python-pptx imports (deferred to give a clean error on missing deps)
# ----------------------------------------------------------------------

def _import_deps():
    try:
        from pptx import Presentation
        from pptx.util import Emu, Pt, Inches
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from lxml import etree
        return Presentation, Emu, Pt, Inches, MSO_SHAPE_TYPE, etree
    except Exception as e:
        print(json.dumps({
            "error": "missing_dependencies",
            "msg": "Install with: python3 -m pip install python-pptx lxml Pillow",
            "detail": str(e),
        }))
        sys.exit(2)

# ----------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------

def get_runs(text_frame):
    """Yield (run, paragraph_index) for every run in the frame."""
    for pi, p in enumerate(text_frame.paragraphs):
        for r in p.runs:
            yield r, pi

def shape_text(shape):
    if not shape.has_text_frame:
        return ''
    return '\n'.join(p.text for p in shape.text_frame.paragraphs)

def char_factor(font_name, bold):
    if not font_name:
        return DEFAULT_CHAR_FACTOR
    return CHAR_WIDTH_FACTORS.get((font_name, bool(bold)), DEFAULT_CHAR_FACTOR)

def estimate_word_width_emu(word, font_name, font_size_pt, bold):
    """EMU width estimate for a single word."""
    factor = char_factor(font_name, bold)
    font_size_emu = font_size_pt * EMU_PER_PT
    return int(len(word) * factor * font_size_emu * WORD_SAFETY)

def simulate_wrap(text, box_w_emu, font_size_pt, font_name, bold):
    """Approximate PowerPoint word-wrap. Returns line count."""
    if not text or not text.strip():
        return 1
    factor = char_factor(font_name, bold)
    font_size_emu = font_size_pt * EMU_PER_PT
    avg_char_w = factor * font_size_emu
    space_w = avg_char_w * 0.35
    usable_w = box_w_emu * 0.95
    lines = 1
    current_w = 0
    for paragraph in text.split('\n'):
        # Each \n forces a new line in the count
        if paragraph != text.split('\n')[0]:
            lines += 1
            current_w = 0
        for word in paragraph.split():
            word_w = len(word) * avg_char_w
            test_w = current_w + (space_w if current_w > 0 else 0) + word_w
            if current_w > 0 and test_w > usable_w:
                lines += 1
                current_w = word_w
            else:
                current_w = test_w
    return max(1, lines)

def first_run_meta(text_frame, default_font='Calibri', default_size=18):
    """Return (font_name, font_size_pt, bold) from first run of frame."""
    for r, _ in get_runs(text_frame):
        size = r.font.size.pt if r.font.size else default_size
        return (r.font.name or default_font, size, bool(r.font.bold))
    return (default_font, default_size, False)

def is_text_box(shape, MSO_SHAPE_TYPE):
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    except Exception:
        return False

def has_image(shape):
    try:
        return shape.shape_type and hasattr(shape, 'image')
    except Exception:
        return False

def shape_bbox(shape):
    """Return (l, t, r, b) tuple in EMU. None if not positioned."""
    try:
        return (shape.left, shape.top,
                shape.left + shape.width, shape.top + shape.height)
    except Exception:
        return None

def bbox_overlap_ratio(a, b):
    """Fraction of the smaller bbox covered by the intersection."""
    al, at, ar, ab = a
    bl, bt, br, bb = b
    ix1, iy1 = max(al, bl), max(at, bt)
    ix2, iy2 = min(ar, br), min(ab, bb)
    iw, ih = max(0, ix2 - ix1), max(0, iy2 - iy1)
    iarea = iw * ih
    if iarea == 0:
        return 0.0
    aarea = max(1, (ar - al) * (ab - at))
    barea = max(1, (br - bl) * (bb - bt))
    return iarea / min(aarea, barea)

def center_in(small, big):
    """Is the center of `small` inside `big`?"""
    sl, st, sr, sb = small
    bl, bt, br, bb = big
    cx, cy = (sl + sr) / 2, (st + sb) / 2
    return bl <= cx <= br and bt <= cy <= bb

def issue(severity, check, slide, shape, msg, **extra):
    out = {
        "severity": severity,
        "check": check,
        "slide": slide,
        "shape": shape,
        "msg": msg,
    }
    out.update(extra)
    return out

# ----------------------------------------------------------------------
# CHECK 1 — BOUNDS
# ----------------------------------------------------------------------

def check_bounds(prs, MSO_SHAPE_TYPE):
    issues = []
    SW, SH = prs.slide_width, prs.slide_height
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            bbox = shape_bbox(shape)
            if bbox is None:
                continue
            l, t, r, b = bbox
            if l < -BOUNDS_TOL:
                issues.append(issue("CRITICAL", 1, si, shape.name,
                    f"Shape left={l} < 0 (off-slide on left)"))
            if t < -BOUNDS_TOL:
                issues.append(issue("CRITICAL", 1, si, shape.name,
                    f"Shape top={t} < 0 (off-slide on top)"))
            if r > SW + BOUNDS_TOL:
                issues.append(issue("CRITICAL", 1, si, shape.name,
                    f"Shape right={r} exceeds slide_width={SW} (off-slide on right)"))
            if b > SH + BOUNDS_TOL:
                issues.append(issue("CRITICAL", 1, si, shape.name,
                    f"Shape bottom={b} exceeds slide_height={SH} (off-slide on bottom)"))
    return issues

# ----------------------------------------------------------------------
# CHECK 2 — TEXT CLIPPING (vertical overflow)
# ----------------------------------------------------------------------

def check_text_clipping(prs, MSO_SHAPE_TYPE):
    issues = []
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape_text(shape)
            if not text.strip():
                continue
            font_name, font_pt, bold = first_run_meta(shape.text_frame)
            try:
                box_w = shape.width
                box_h = shape.height
            except Exception:
                continue
            if box_w <= 0 or box_h <= 0:
                continue
            lines = simulate_wrap(text, box_w, font_pt, font_name, bold)
            est_h = lines * font_pt * EMU_PER_PT * LINE_HEIGHT_RATIO
            if est_h > box_h * TEXT_OVERFLOW_TOL:
                issues.append(issue("WARNING", 2, si, shape.name,
                    f"Estimated text height {est_h/EMU_PER_PT:.1f}pt > frame height "
                    f"{box_h/EMU_PER_PT:.1f}pt (lines={lines}, font={font_pt}pt). "
                    f"Estimator may overshoot ~30%; visually verify before fixing.",
                    lines=lines, font_pt=font_pt))
    return issues

# ----------------------------------------------------------------------
# CHECK 3 — WORD-WRAP QUALITY (longest word fit)
# ----------------------------------------------------------------------

def check_word_wrap(prs, MSO_SHAPE_TYPE):
    issues = []
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            try:
                box_w = shape.width
            except Exception:
                continue
            if box_w <= 0:
                continue
            for run, _ in get_runs(shape.text_frame):
                if not run.text:
                    continue
                size_pt = run.font.size.pt if run.font.size else 18
                font_name = run.font.name or 'Calibri'
                bold = bool(run.font.bold)
                for word in run.text.split():
                    if len(word) <= 4:
                        continue
                    w_emu = estimate_word_width_emu(word, font_name, size_pt, bold)
                    if w_emu > box_w:
                        issues.append(issue("CRITICAL", 3, si, shape.name,
                            f"Long word '{word}' ({len(word)} chars at {size_pt}pt {font_name}) "
                            f"width {w_emu/EMU_PER_INCH:.2f}\" > frame width "
                            f"{box_w/EMU_PER_INCH:.2f}\" — will overflow",
                            word=word, size_pt=size_pt))
                        break  # one critical per shape is enough
                    elif w_emu > box_w * 0.96 and len(word) > 5:
                        issues.append(issue("WARNING", 3, si, shape.name,
                            f"Tight word '{word}' at {w_emu/box_w:.0%} of frame width",
                            word=word, size_pt=size_pt))
                        break
    return issues

# ----------------------------------------------------------------------
# CHECK 4 — CONTAINER / TEXT SYNC
# ----------------------------------------------------------------------

def check_container_text_sync(prs, MSO_SHAPE_TYPE):
    issues = []
    PADDING = int(0.04 * EMU_PER_INCH)
    for si, slide in enumerate(prs.slides, 1):
        shapes = list(slide.shapes)
        # Find candidate containers: non-image, non-textbox shapes with positive area >= ~1"x1"
        containers = []
        text_boxes = []
        for sh in shapes:
            bbox = shape_bbox(sh)
            if bbox is None:
                continue
            l, t, r, b = bbox
            w, h = r - l, b - t
            if w < EMU_PER_INCH or h < EMU_PER_INCH * 0.5:
                continue
            if has_image(sh):
                continue
            if is_text_box(sh, MSO_SHAPE_TYPE):
                text_boxes.append(sh)
            else:
                containers.append(sh)
        for tb in text_boxes:
            tb_bbox = shape_bbox(tb)
            if tb_bbox is None:
                continue
            for cont in containers:
                cb = shape_bbox(cont)
                if cb is None:
                    continue
                if not center_in(tb_bbox, cb):
                    continue
                cl, ct, cr, cb_ = cb
                tl, tt, tr, tb_ = tb_bbox
                if (tl < cl + PADDING - BOUNDS_TOL or
                    tt < ct + PADDING - BOUNDS_TOL or
                    tr > cr - PADDING + BOUNDS_TOL or
                    tb_ > cb_ - PADDING + BOUNDS_TOL):
                    issues.append(issue("CRITICAL", 4, si, tb.name,
                        f"Text box '{tb.name}' overflows its container '{cont.name}' "
                        f"(text bbox not inside container - padding). "
                        f"Grow container or shrink text/font.",
                        container=cont.name))
                    break
    return issues

# ----------------------------------------------------------------------
# CHECK 5 — BULLET ALIGNMENT (conservative)
# ----------------------------------------------------------------------

def check_bullet_alignment(prs, MSO_SHAPE_TYPE):
    """Conservative: only flag groups of >= 3 small-shape "dots" whose
    horizontal positions vary by more than 5000 EMU. Strict line-count
    layout check (paired-text per dot) is left to the agent — too many
    false positives from auto-detection alone."""
    issues = []
    DOT_MAX = int(0.15 * EMU_PER_INCH)
    for si, slide in enumerate(prs.slides, 1):
        dots = []
        for sh in slide.shapes:
            if has_image(sh):
                continue
            if is_text_box(sh, MSO_SHAPE_TYPE):
                continue
            try:
                if sh.width <= DOT_MAX and sh.height <= DOT_MAX:
                    if not shape_text(sh).strip():
                        dots.append(sh)
            except Exception:
                continue
        # Group dots that share a vertical column (likely a single bullet list)
        cols = defaultdict(list)
        for d in dots:
            cols[d.left // 50000].append(d)  # bucket by ~0.05"
        for col_dots in cols.values():
            if len(col_dots) < 3:
                continue
            lefts = [d.left for d in col_dots]
            if max(lefts) - min(lefts) > 5000:
                issues.append(issue("WARNING", 5, si, ",".join(d.name for d in col_dots[:5]),
                    f"Bullet dots in same list have inconsistent left positions "
                    f"(spread = {(max(lefts)-min(lefts))/EMU_PER_INCH:.3f}\"). "
                    f"Re-align to a single column."))
    return issues

# ----------------------------------------------------------------------
# CHECK 6 — OVERLAP CLASSIFICATION
# ----------------------------------------------------------------------

def check_overlap(prs, MSO_SHAPE_TYPE):
    """Flag pairs of independent text shapes whose bboxes overlap > 10%.
    Skip parent-child (one center inside the other) and full-slide image/overlay stacks."""
    issues = []
    SW, SH = prs.slide_width, prs.slide_height
    for si, slide in enumerate(prs.slides, 1):
        text_shapes = []
        for sh in slide.shapes:
            if not shape_text(sh).strip():
                continue
            bbox = shape_bbox(sh)
            if bbox is None:
                continue
            text_shapes.append((sh, bbox))
        n = len(text_shapes)
        for i in range(n):
            sa, ba = text_shapes[i]
            for j in range(i + 1, n):
                sb, bb = text_shapes[j]
                # skip if either is full-slide
                def is_fullslide(b):
                    l, t, r, btm = b
                    return (r - l) > SW * 0.95 and (btm - t) > SH * 0.95
                if is_fullslide(ba) or is_fullslide(bb):
                    continue
                # skip if one center is inside the other (parent-child)
                if center_in(ba, bb) or center_in(bb, ba):
                    continue
                ratio = bbox_overlap_ratio(ba, bb)
                if ratio > 0.10:
                    issues.append(issue("WARNING", 6, si,
                        f"{sa.name} ↔ {sb.name}",
                        f"Two independent text shapes overlap by {ratio:.0%}",
                        ratio=round(ratio, 2)))
    return issues

# ----------------------------------------------------------------------
# CHECK 7 — Z-ORDER (conservative)
# ----------------------------------------------------------------------

def check_z_order(prs, MSO_SHAPE_TYPE):
    """Detect opaque-fill shape layered above (later in shape list = above)
    a text shape, where the cover fraction of the *text shape's area* is
    > 30%. Conservative: skip everything that doesn't have an obviously
    opaque solid fill, and skip thin accent bars (< 5pt in either
    dimension) which routinely sit inside title boxes as underlines."""
    issues = []
    THIN_EMU = 5 * EMU_PER_PT  # ~5pt
    for si, slide in enumerate(prs.slides, 1):
        shapes = list(slide.shapes)
        for i, sh in enumerate(shapes):
            if has_image(sh) or is_text_box(sh, MSO_SHAPE_TYPE):
                continue
            try:
                fill = sh.fill
                if fill.type != 1:  # MSO_FILL.SOLID
                    continue
            except Exception:
                continue
            sh_bbox = shape_bbox(sh)
            if sh_bbox is None:
                continue
            # Skip thin accent bars
            try:
                if sh.height < THIN_EMU or sh.width < THIN_EMU:
                    continue
            except Exception:
                continue
            sh_area = (sh_bbox[2] - sh_bbox[0]) * (sh_bbox[3] - sh_bbox[1])
            for j in range(i):
                other = shapes[j]
                if not shape_text(other).strip():
                    continue
                ob = shape_bbox(other)
                if ob is None:
                    continue
                # Intersection area
                ix1, iy1 = max(sh_bbox[0], ob[0]), max(sh_bbox[1], ob[1])
                ix2, iy2 = min(sh_bbox[2], ob[2]), min(sh_bbox[3], ob[3])
                iw, ih = max(0, ix2 - ix1), max(0, iy2 - iy1)
                iarea = iw * ih
                if iarea == 0:
                    continue
                text_area = (ob[2] - ob[0]) * (ob[3] - ob[1])
                if text_area == 0:
                    continue
                cover_of_text = iarea / text_area
                if cover_of_text > 0.30:
                    issues.append(issue("CRITICAL", 7, si, sh.name,
                        f"Opaque fill shape '{sh.name}' is z-ordered above text shape "
                        f"'{other.name}' and covers {cover_of_text:.0%} of the text area",
                        covered_by=sh.name, covered_text=other.name,
                        cover_of_text_pct=round(cover_of_text * 100, 1)))
    return issues

# ----------------------------------------------------------------------
# CHECK 8 — FONT COMPLIANCE
# ----------------------------------------------------------------------

def check_font_compliance(prs, MSO_SHAPE_TYPE):
    issues = []
    for si, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            for r, _ in get_runs(sh.text_frame):
                if not r.text or not r.text.strip():
                    continue
                size = r.font.size
                name = r.font.name
                if size is None:
                    issues.append(issue("WARNING", 8, si, sh.name,
                        f"Run '{r.text[:40]}' has no explicit font.size — theme defaults are unreliable"))
                else:
                    if size.pt < MIN_FONT_PT:
                        sev = "WARNING" if size.pt >= MIN_FONT_PT_CAPTION else "CRITICAL"
                        issues.append(issue(sev, 8, si, sh.name,
                            f"Run '{r.text[:40]}' is {size.pt}pt — below {MIN_FONT_PT}pt floor",
                            size_pt=size.pt))
                if name is None:
                    issues.append(issue("WARNING", 8, si, sh.name,
                        f"Run '{r.text[:40]}' has no explicit font.name — theme defaults are unreliable"))
    return issues

# ----------------------------------------------------------------------
# CHECK 9 — SPACING CONSISTENCY (cross-slide left margins)
# ----------------------------------------------------------------------

def check_spacing(prs, MSO_SHAPE_TYPE):
    """Conservative cross-slide left-margin check. We look at the
    leftmost text shape on each slide that's wider than 30% of slide
    width (so titles, body blocks) and flag if their lefts vary by
    more than 100,000 EMU (~0.11")."""
    issues = []
    SW = prs.slide_width
    primary_lefts = []
    for si, slide in enumerate(prs.slides, 1):
        candidates = []
        for sh in slide.shapes:
            if not shape_text(sh).strip():
                continue
            try:
                if sh.width >= SW * 0.30:
                    candidates.append((sh.left, sh.name, si))
            except Exception:
                continue
        if candidates:
            candidates.sort()
            primary_lefts.append(candidates[0])
    if len(primary_lefts) >= 3:
        lefts = [pl[0] for pl in primary_lefts]
        spread = max(lefts) - min(lefts)
        if spread > 100000:
            issues.append(issue("INFO", 9, 0, "deck-wide",
                f"Primary left-margin spread across slides = "
                f"{spread/EMU_PER_INCH:.3f}\" (>0.11\"). Verify titles "
                f"are aligned consistently.",
                spread_emu=spread))
    return issues

# ----------------------------------------------------------------------
# CHECK 10 — COLOR / FILL INTEGRITY (orphan accent1 fills)
# ----------------------------------------------------------------------

def check_color_integrity(prs, MSO_SHAPE_TYPE, etree):
    """Flag medium-or-larger shapes carrying a theme accent1 reference
    with no explicit override. These render as a flat blue rectangle in
    PowerPoint when the author intended a custom color/gradient."""
    issues = []
    SW, SH = prs.slide_width, prs.slide_height
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    for si, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if has_image(sh) or is_text_box(sh, MSO_SHAPE_TYPE):
                continue
            sp = sh.element
            spPr = sp.find(f'{{{ns_p}}}spPr')
            p_style = sp.find(f'{{{ns_p}}}style')
            if spPr is None or p_style is None:
                continue
            has_explicit = (
                spPr.find(f'{{{ns_a}}}solidFill') is not None or
                spPr.find(f'{{{ns_a}}}gradFill') is not None or
                spPr.find(f'{{{ns_a}}}noFill') is not None or
                spPr.find(f'{{{ns_a}}}pattFill') is not None
            )
            try:
                style_xml = etree.tostring(p_style).decode()
            except Exception:
                continue
            if 'accent1' in style_xml and not has_explicit:
                try:
                    is_med = (sh.width > SW * 0.20 or sh.height > SH * 0.15)
                except Exception:
                    is_med = False
                if is_med:
                    issues.append(issue("WARNING", 10, si, sh.name,
                        f"'{sh.name}' uses theme accent1 with no explicit fill override "
                        f"— may render as theme blue instead of intended color"))
    return issues

# ----------------------------------------------------------------------
# CHECK 11 — STYLE COMPLIANCE (only when --style is passed)
# ----------------------------------------------------------------------

# Minimal embedded style dict subset. The skill's full mapping lives in
# pptx-design/references/style-pptx-mapping.md. We only encode the parts
# the auditor can check programmatically: required font names + palette
# RGB values. If a style isn't here, CHECK 11 returns no issues.
_STYLE_DICTS = {
    "STYLE-01": {  # Strategy Consulting
        "fonts": ["Georgia", "Calibri", "Arial"],
        "palette": [(31, 58, 88), (255, 255, 255), (212, 175, 55), (180, 35, 25), (90, 90, 90)],
    },
    "STYLE-02": {  # Executive Editorial
        "fonts": ["Georgia", "Calibri"],
        "palette": [(20, 30, 48), (250, 245, 235), (200, 70, 60), (140, 110, 70)],
    },
    "STYLE-03": {  # Sketch / Hand-Drawn
        "fonts": ["Comic Sans MS", "Calibri"],
        "palette": [(255, 250, 240), (40, 40, 40), (255, 130, 100)],
    },
    # Other styles intentionally omitted for v1 — add to the dict if needed.
}

def check_style_compliance(prs, style_id, MSO_SHAPE_TYPE):
    issues = []
    spec = _STYLE_DICTS.get(style_id)
    if spec is None:
        issues.append(issue("INFO", 11, 0, "deck-wide",
            f"Style '{style_id}' not in audit's known style dict — CHECK 11 skipped. "
            f"Add an entry to _STYLE_DICTS in audit.py to enforce."))
        return issues
    expected_fonts = set(spec["fonts"])
    expected_palette = set(spec["palette"])
    found_fonts = set()
    found_colors = set()
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                for r, _ in get_runs(sh.text_frame):
                    if r.font.name:
                        found_fonts.add(r.font.name)
            try:
                if sh.fill.type == 1:
                    rgb = sh.fill.fore_color.rgb
                    if rgb is not None:
                        found_colors.add((rgb[0], rgb[1], rgb[2]))
            except Exception:
                pass
    extra_fonts = found_fonts - expected_fonts
    if extra_fonts:
        issues.append(issue("WARNING", 11, 0, "deck-wide",
            f"Off-style fonts: {sorted(extra_fonts)} (style {style_id} expects {sorted(expected_fonts)})",
            offending_fonts=sorted(extra_fonts)))
    # Color check tolerance: ±10 per channel
    def near_palette(c):
        for ec in expected_palette:
            if all(abs(c[i] - ec[i]) <= 10 for i in range(3)):
                return True
        return False
    extra_colors = [c for c in found_colors
                    if c != (255, 255, 255) and c != (0, 0, 0)
                    and not near_palette(c)]
    if extra_colors:
        issues.append(issue("WARNING", 11, 0, "deck-wide",
            f"Off-palette colors: {extra_colors} (style {style_id} palette: {sorted(expected_palette)})",
            offending_colors=extra_colors))
    return issues

# ----------------------------------------------------------------------
# CHECK 12 — IMAGE ASPECT RATIO DISTORTION
# ----------------------------------------------------------------------

def check_image_ar(prs, MSO_SHAPE_TYPE):
    issues = []
    try:
        from PIL import Image as PILImage
    except Exception:
        issues.append(issue("INFO", 12, 0, "deck-wide",
            "Pillow not installed — CHECK 12 (image AR distortion) skipped. "
            "Install with: python3 -m pip install Pillow"))
        return issues
    SW, SH = prs.slide_width, prs.slide_height
    for si, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not has_image(sh):
                continue
            try:
                blob = sh.image.blob
                img = PILImage.open(io.BytesIO(blob))
                native_w, native_h = img.size
                img.close()
            except Exception:
                continue
            try:
                placed_w, placed_h = sh.width, sh.height
            except Exception:
                continue
            if native_h == 0 or placed_h == 0:
                continue
            native_ar = native_w / native_h
            placed_ar = placed_w / placed_h
            # skip full-bleed BG (always 16:9 → 16:9 in 16:9 deck)
            try:
                if (abs(sh.left) < 50000 and abs(sh.top) < 50000 and
                    abs(placed_w - SW) < 50000 and abs(placed_h - SH) < 50000):
                    continue
            except Exception:
                pass
            distortion = abs(native_ar - placed_ar) / native_ar
            if distortion > AR_CRITICAL:
                issues.append(issue("CRITICAL", 12, si, sh.name,
                    f"Image '{sh.name}': native AR={native_ar:.2f}, "
                    f"placed AR={placed_ar:.2f}, distortion={distortion:.0%}. "
                    f"Use add_picture_fit() or regenerate at correct AR.",
                    native_ar=round(native_ar, 2),
                    placed_ar=round(placed_ar, 2),
                    distortion_pct=round(distortion * 100, 1)))
            elif distortion > AR_WARNING:
                issues.append(issue("WARNING", 12, si, sh.name,
                    f"Image '{sh.name}': minor AR distortion {distortion:.0%}",
                    distortion_pct=round(distortion * 100, 1)))
    return issues

# ----------------------------------------------------------------------
# CHECK 13 — BROKEN GRADIENT FILLS ("blue rectangle" bug)
# ----------------------------------------------------------------------

def check_broken_gradients(prs, MSO_SHAPE_TYPE, etree):
    issues = []
    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    SW, SH = prs.slide_width, prs.slide_height
    for si, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if has_image(sh) or is_text_box(sh, MSO_SHAPE_TYPE):
                continue
            sp = sh.element
            stray = sp.findall(f'{{{ns_a}}}gradFill')
            if stray:
                issues.append(issue("CRITICAL", 13, si, sh.name,
                    f"'{sh.name}': gradFill attached to <p:sp> instead of <p:spPr> — "
                    f"PowerPoint will ignore the gradient and render theme blue. "
                    f"Use add_gradient_shape() helper."))
                continue
            spPr = sp.find(f'{{{ns_p}}}spPr')
            p_style = sp.find(f'{{{ns_p}}}style')
            if spPr is None or p_style is None:
                continue
            has_explicit = (
                spPr.find(f'{{{ns_a}}}solidFill') is not None or
                spPr.find(f'{{{ns_a}}}gradFill') is not None or
                spPr.find(f'{{{ns_a}}}noFill') is not None or
                spPr.find(f'{{{ns_a}}}pattFill') is not None
            )
            try:
                style_xml = etree.tostring(p_style).decode()
            except Exception:
                continue
            if 'accent1' in style_xml and not has_explicit:
                try:
                    is_large = (sh.width > SW * 0.5 or sh.height > SH * 0.3)
                except Exception:
                    is_large = False
                if is_large:
                    issues.append(issue("CRITICAL", 13, si, sh.name,
                        f"'{sh.name}': large shape with theme accent1 fill, "
                        f"no explicit spPr fill — likely the 'blue rectangle' bug"))
    return issues

# ----------------------------------------------------------------------
# CHECK 14 — TEXT ZONE LUMINANCE
# ----------------------------------------------------------------------

_ANNOT_RE = re.compile(r'text_zone=(?P<zone>\w+):(?P<size>[0-9.]+):(?P<color>\w+)')

def check_text_zone_luminance(prs, MSO_SHAPE_TYPE):
    issues = []
    try:
        from PIL import Image as PILImage
    except Exception:
        return issues  # silently skip — already reported in CHECK 12
    for si, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            if not has_image(sh):
                continue
            m = _ANNOT_RE.search(sh.name or '')
            if not m:
                continue
            zone = m.group('zone')
            try:
                size = float(m.group('size'))
            except Exception:
                continue
            color = m.group('color').lower()
            try:
                img = PILImage.open(io.BytesIO(sh.image.blob)).convert('L')
            except Exception:
                continue
            W, H = img.size
            # Crop the text zone
            if zone == 'bottom':
                box = (0, int(H * (1 - size)), W, H)
            elif zone == 'top':
                box = (0, 0, W, int(H * size))
            elif zone == 'left':
                box = (0, 0, int(W * size), H)
            elif zone == 'right':
                box = (int(W * (1 - size)), 0, W, H)
            else:
                continue
            crop = img.crop(box)
            # Mean luminance
            px = list(crop.getdata())
            mean_lum = sum(px) / max(1, len(px))
            light_text = color in ('white', 'cream', 'lightcream', 'ivory')
            dark_text = color in ('black', 'dark', 'navy')
            if light_text and mean_lum > LUM_CRITICAL_BRIGHT:
                issues.append(issue("CRITICAL", 14, si, sh.name,
                    f"Text zone '{zone}' on image '{sh.name}' has mean luminance "
                    f"{mean_lum:.0f} (> {LUM_CRITICAL_BRIGHT}) — too bright for "
                    f"{color} text. Add gradient overlay, regen image with darker zone, or flip color.",
                    mean_luminance=round(mean_lum, 1), text_color=color))
            elif dark_text and mean_lum < LUM_CRITICAL_DARK:
                issues.append(issue("CRITICAL", 14, si, sh.name,
                    f"Text zone '{zone}' on image '{sh.name}' has mean luminance "
                    f"{mean_lum:.0f} (< {LUM_CRITICAL_DARK}) — too dark for "
                    f"{color} text.",
                    mean_luminance=round(mean_lum, 1), text_color=color))
    return issues

# ----------------------------------------------------------------------
# Runner
# ----------------------------------------------------------------------

def run_audit(path, style=None):
    Presentation, Emu, Pt, Inches, MSO_SHAPE_TYPE, etree = _import_deps()
    if not os.path.isfile(path):
        return {
            "path": path,
            "error": "file_not_found",
            "msg": f"File does not exist: {path}",
            "summary": {"critical": 1, "warning": 0, "info": 0, "passed": False},
            "critical": [{"check": 0, "slide": 0, "shape": "", "msg": "file not found"}],
            "warning": [], "info": [],
        }
    prs = Presentation(path)
    issues = []
    issues.extend(check_bounds(prs, MSO_SHAPE_TYPE))
    issues.extend(check_text_clipping(prs, MSO_SHAPE_TYPE))
    issues.extend(check_word_wrap(prs, MSO_SHAPE_TYPE))
    issues.extend(check_container_text_sync(prs, MSO_SHAPE_TYPE))
    issues.extend(check_bullet_alignment(prs, MSO_SHAPE_TYPE))
    issues.extend(check_overlap(prs, MSO_SHAPE_TYPE))
    issues.extend(check_z_order(prs, MSO_SHAPE_TYPE))
    issues.extend(check_font_compliance(prs, MSO_SHAPE_TYPE))
    issues.extend(check_spacing(prs, MSO_SHAPE_TYPE))
    issues.extend(check_color_integrity(prs, MSO_SHAPE_TYPE, etree))
    if style:
        issues.extend(check_style_compliance(prs, style, MSO_SHAPE_TYPE))
    issues.extend(check_image_ar(prs, MSO_SHAPE_TYPE))
    issues.extend(check_broken_gradients(prs, MSO_SHAPE_TYPE, etree))
    issues.extend(check_text_zone_luminance(prs, MSO_SHAPE_TYPE))
    sev = lambda s: [i for i in issues if i["severity"] == s]
    crit = sev("CRITICAL")
    warn = sev("WARNING")
    info = sev("INFO")
    return {
        "path": os.path.abspath(path),
        "style": style,
        "summary": {
            "critical": len(crit),
            "warning": len(warn),
            "info": len(info),
            "passed": len(crit) == 0,
        },
        "critical": crit,
        "warning": warn,
        "info": info,
    }

def main():
    ap = argparse.ArgumentParser(
        description="Deterministic 14-check audit for PowerPoint decks. "
                    "Exit 0 if no CRITICAL issues, 1 otherwise.")
    ap.add_argument("path", help="Path to .pptx file")
    ap.add_argument("--style", default=None,
                    help="Style ID for CHECK 11 (e.g. STYLE-02). Skipped if omitted.")
    ap.add_argument("--fail-on", choices=["critical", "warning"], default="critical",
                    help="Exit nonzero on this severity or higher (default: critical).")
    ap.add_argument("--json-only", action="store_true",
                    help="Suppress the human-readable stderr summary.")
    args = ap.parse_args()
    report = run_audit(args.path, style=args.style)
    print(json.dumps(report, indent=2))
    if not args.json_only:
        s = report.get("summary", {})
        msg = (f"\n[pptx-audit] CRITICAL={s.get('critical',0)} "
               f"WARNING={s.get('warning',0)} INFO={s.get('info',0)} "
               f"passed={s.get('passed', False)}")
        print(msg, file=sys.stderr)
    crit = report.get("summary", {}).get("critical", 0)
    warn = report.get("summary", {}).get("warning", 0)
    if args.fail_on == "critical":
        sys.exit(0 if crit == 0 else 1)
    return sys.exit(0 if (crit + warn) == 0 else 1)

if __name__ == "__main__":
    main()
