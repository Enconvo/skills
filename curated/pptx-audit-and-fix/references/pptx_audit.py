#!/usr/bin/env python3
"""
PPTX Layout Audit & Fix Tool
Works outside of PowerPoint — operates directly on .pptx files using python-pptx.
Dependencies: pip install python-pptx Pillow lxml
Usage:
    from pptx_audit import PptxAuditor
    auditor = PptxAuditor("presentation.pptx")
    report = auditor.run_full_audit()
    auditor.fix_all(report)
    auditor.save("presentation_fixed.pptx")
"""
import copy, math, os, re
from dataclasses import dataclass, field
from enum import Enum
from typing import Optional
from lxml import etree
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}
MIN_MARGIN_PT = 36   # 0.5 inch
MIN_FONT_SIZE = 14
MIN_BODY_FONT_SIZE = 16
MIN_OVERFLOW_FIX_PT = 8  # Only auto-fix text overflows > this threshold (PowerPoint handles small ones)
CHAR_WIDTH_RATIO = {
    "default": 0.52, "Georgia": 0.55, "Arial": 0.50, "Calibri": 0.49,
    "Consolas": 0.60, "Courier New": 0.60, "Times New Roman": 0.51,
    "Rockwell": 0.56, "Helvetica": 0.50,
}
LINE_HEIGHT_RATIO = 1.25


# ── Data Classes ──────────────────────────────────────────────────────────────

class Severity(Enum):
    CRITICAL = "CRITICAL"
    WARNING = "WARNING"
    INFO = "INFO"

@dataclass
class BBox:
    left: float; top: float; width: float; height: float
    @property
    def right(self): return self.left + self.width
    @property
    def bottom(self): return self.top + self.height

@dataclass
class ShapeInfo:
    shape_id: str; name: str; shape_type: str; bbox: BBox
    text: str = ""; font_size_pt: float = 0; font_name: str = ""
    font_color: str = ""; font_bold: bool = False; font_italic: bool = False
    is_background: bool = False; is_decorative: bool = False
    is_page_furniture: bool = False; has_fill: bool = False; estimated_text_height: float = 0
    has_text_overflow: bool = False
    _shape_ref: object = field(default=None, repr=False)

@dataclass
class Issue:
    slide_num: int; severity: Severity; category: str; description: str
    shape_a: Optional[str] = None; shape_b: Optional[str] = None
    fix_description: str = ""; auto_fixable: bool = False

@dataclass
class AuditReport:
    filename: str; total_slides: int
    slide_width_pt: float; slide_height_pt: float
    issues: list = field(default_factory=list)
    @property
    def critical_count(self): return sum(1 for i in self.issues if i.severity == Severity.CRITICAL)
    @property
    def warning_count(self): return sum(1 for i in self.issues if i.severity == Severity.WARNING)
    @property
    def info_count(self): return sum(1 for i in self.issues if i.severity == Severity.INFO)
    def __str__(self):
        lines = [f"PPTX Audit Report — {self.filename}", "=" * 50,
                 f"Slides: {self.total_slides}, Critical: {self.critical_count}, "
                 f"Warnings: {self.warning_count}, Info: {self.info_count}", ""]
        from itertools import groupby
        sorted_issues = sorted(self.issues, key=lambda i: (i.slide_num, i.severity.value))
        for slide_num, group in groupby(sorted_issues, key=lambda i: i.slide_num):
            lines.append(f"\nSlide {slide_num}:")
            lines.append("-" * 40)
            for issue in group:
                icon = {"CRITICAL": "🔴", "WARNING": "🟡", "INFO": "🟢"}[issue.severity.value]
                lines.append(f"  {icon} [{issue.category}] {issue.description}")
                if issue.fix_description:
                    lines.append(f"     → Fix: {issue.fix_description}")
        return "\n".join(lines)


# ── Text Height Estimation ────────────────────────────────────────────────────

def estimate_text_height(text: str, font_size_pt: float, font_name: str,
                         box_width_pt: float, bold: bool = False) -> float:
    """
    Estimate rendered height of text in a box.
    This is the Claude Code replacement for PowerPoint's
    autoSizeSetting = "AutoSizeShapeToFitText" trick.
    Since we can't ask PowerPoint to render, we estimate using:
    1. Character width ratio (empirical per-font)
    2. Word wrapping simulation
    3. Line height calculation
    Returns estimated height in points.
    """
    if not text or font_size_pt <= 0 or box_width_pt <= 0:
        return 0
    ratio = CHAR_WIDTH_RATIO.get(font_name, CHAR_WIDTH_RATIO["default"])
    if bold:
        ratio *= 1.08  # Bold is ~8% wider
    avg_char_width = font_size_pt * ratio
    line_height = font_size_pt * LINE_HEIGHT_RATIO
    paragraphs = text.split("\n")
    total_lines = 0
    for para in paragraphs:
        if not para.strip():
            total_lines += 1
            continue
        words = para.split()
        current_line_width = 0
        lines_in_para = 1
        for word in words:
            word_width = len(word) * avg_char_width
            space_width = avg_char_width
            if current_line_width + word_width > box_width_pt and current_line_width > 0:
                lines_in_para += 1
                current_line_width = word_width + space_width
            else:
                current_line_width += word_width + space_width
        total_lines += lines_in_para
    estimated_height = (total_lines * line_height) + 8  # +8 for padding
    return estimated_height

def emu_to_pt(emu):
    if emu is None: return 0
    return emu / 12700


# ── Shape Extraction ──────────────────────────────────────────────────────────

def extract_shape_info(shape, slide_width_pt, slide_height_pt) -> ShapeInfo:
    bbox = BBox(
        left=emu_to_pt(shape.left) if shape.left is not None else 0,
        top=emu_to_pt(shape.top) if shape.top is not None else 0,
        width=emu_to_pt(shape.width) if shape.width is not None else 0,
        height=emu_to_pt(shape.height) if shape.height is not None else 0,
    )
    shape_type = "other"
    from pptx.shapes.graphfrm import GraphicFrame
    from pptx.shapes.group import GroupShape
    if hasattr(shape, "image"): shape_type = "image"
    elif hasattr(shape, "text_frame"):
        if shape.height and shape.height < Pt(5): shape_type = "geometric"
        else: shape_type = "textbox"
    elif isinstance(shape, GraphicFrame):
        if shape.has_chart: shape_type = "chart"
        elif shape.has_table: shape_type = "table"
    elif isinstance(shape, GroupShape): shape_type = "group"
    else: shape_type = "geometric"
    text = ""; font_size_pt = 0; font_name = ""; font_color = ""
    font_bold = False; font_italic = False
    if shape.has_text_frame:
        tf = shape.text_frame
        text = tf.text or ""
        # Read font properties from raw XML to avoid python-pptx mutating the XML
        # (accessing run.font.* causes python-pptx to inject empty <a:solidFill/> and <a:rPr/> elements)
        from lxml import etree
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        for para_el in shape._element.findall('.//a:p', nsmap):
            for r_el in para_el.findall('a:r', nsmap):
                t_el = r_el.find('a:t', nsmap)
                if t_el is not None and t_el.text and t_el.text.strip():
                    rPr = r_el.find('a:rPr', nsmap)
                    if rPr is not None:
                        sz = rPr.get('sz')
                        if sz: font_size_pt = int(sz) / 100  # hundredths of a point
                        lang_or_typeface = rPr.find('a:latin', nsmap)
                        if lang_or_typeface is not None:
                            tf_name = lang_or_typeface.get('typeface')
                            if tf_name: font_name = tf_name
                        if rPr.get('b') == '1': font_bold = True
                        if rPr.get('i') == '1': font_italic = True
                        srgb = rPr.find('.//a:solidFill/a:srgbClr', nsmap)
                        if srgb is not None:
                            font_color = srgb.get('val', '')
                    break
            if font_size_pt > 0: break
        # Fallback: check paragraph-level defRPr if no run-level font size found
        if font_size_pt == 0:
            for para_el in shape._element.findall('.//a:p', nsmap):
                pPr = para_el.find('a:pPr', nsmap)
                if pPr is not None:
                    defRPr = pPr.find('a:defRPr', nsmap)
                    if defRPr is not None:
                        sz = defRPr.get('sz')
                        if sz:
                            font_size_pt = int(sz) / 100
                            # Also grab font name/bold from defRPr if not already set
                            if not font_name:
                                latin = defRPr.find('a:latin', nsmap)
                                if latin is not None:
                                    tf_name = latin.get('typeface')
                                    if tf_name: font_name = tf_name
                            if not font_bold and defRPr.get('b') == '1':
                                font_bold = True
                            break
                if font_size_pt > 0: break
    is_background = (shape_type == "image" and bbox.width >= slide_width_pt * 0.9
                     and bbox.height >= slide_height_pt * 0.9)
    is_decorative = (shape_type == "geometric" and bbox.height <= 5)
    # Check if shape has a visible fill — read from raw XML to avoid python-pptx side effects
    # Fills can be in spPr (explicit) or via fillRef in style (theme-referenced)
    has_fill = False
    try:
        sp_el = shape._element
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        # Check explicit fill in spPr
        spPr = sp_el.find('.//a:spPr', nsmap) if sp_el is not None else None
        if spPr is not None:
            has_fill = (spPr.find('a:solidFill', nsmap) is not None
                       or spPr.find('a:gradFill', nsmap) is not None
                       or spPr.find('a:pattFill', nsmap) is not None
                       or spPr.find('a:blipFill', nsmap) is not None)
        # Check style-referenced fill (fillRef with idx > 0 means theme fill)
        if not has_fill:
            fillRef = sp_el.find('.//a:fillRef', nsmap)
            if fillRef is not None:
                idx = fillRef.get('idx', '0')
                if idx != '0':  # idx=0 means no fill
                    has_fill = True
    except Exception:
        pass
    is_page_furniture = (bbox.left > slide_width_pt * 0.85
                         and bbox.top > slide_height_pt * 0.85
                         and bbox.width < 100 and bbox.height < 40)
    estimated_height = 0; has_overflow = False
    if text and font_size_pt > 0:
        effective_width = max(bbox.width - 14, 10)
        estimated_height = estimate_text_height(text, font_size_pt, font_name,
                                                 effective_width, font_bold)
        has_overflow = estimated_height > bbox.height * 1.05
    return ShapeInfo(shape_id=str(shape.shape_id), name=shape.name or f"Shape_{shape.shape_id}",
        shape_type=shape_type, bbox=bbox, text=text, font_size_pt=font_size_pt,
        font_name=font_name, font_color=font_color, font_bold=font_bold,
        font_italic=font_italic, is_background=is_background, is_decorative=is_decorative,
        is_page_furniture=is_page_furniture, has_fill=has_fill, estimated_text_height=estimated_height,
        has_text_overflow=has_overflow, _shape_ref=shape)


# ── Overlap Detection + Phase 1: Structural Scan ─────────────────────────────

def calc_overlap(a: BBox, b: BBox) -> tuple:
    overlap_x = min(a.right, b.right) - max(a.left, b.left)
    overlap_y = min(a.bottom, b.bottom) - max(a.top, b.top)
    return (max(0, overlap_x), max(0, overlap_y))

def calc_overlap_with_estimated_height(shape_a: ShapeInfo, shape_b: ShapeInfo) -> tuple:
    a_h = max(shape_a.bbox.height, shape_a.estimated_text_height) if shape_a.estimated_text_height > 0 else shape_a.bbox.height
    b_h = max(shape_b.bbox.height, shape_b.estimated_text_height) if shape_b.estimated_text_height > 0 else shape_b.bbox.height
    a_bottom = shape_a.bbox.top + a_h
    b_bottom = shape_b.bbox.top + b_h
    ox = min(shape_a.bbox.right, shape_b.bbox.right) - max(shape_a.bbox.left, shape_b.bbox.left)
    oy = min(a_bottom, b_bottom) - max(shape_a.bbox.top, shape_b.bbox.top)
    return (max(0, ox), max(0, oy))

def phase1_structural_scan(shapes, slide_num, slide_w, slide_h):
    issues = []
    content_shapes = [s for s in shapes if not s.is_background]
    for i, a in enumerate(content_shapes):
        for b in content_shapes[i + 1:]:
            if "scrim" in a.name.lower() or "overlay" in a.name.lower(): continue
            if "scrim" in b.name.lower() or "overlay" in b.name.lower(): continue
            ox, oy = calc_overlap(a.bbox, b.bbox)
            if ox > 0 and oy > 2:
                if a.is_decorative and not b.is_decorative:
                    issues.append(Issue(slide_num, Severity.CRITICAL, "overlap",
                        f"Decorative line '{a.name}' cuts through '{b.name}' by {oy:.1f}pt",
                        a.name, b.name, f"Move '{a.name}' above or below '{b.name}'", True))
                elif b.is_decorative and not a.is_decorative:
                    issues.append(Issue(slide_num, Severity.CRITICAL, "overlap",
                        f"Decorative line '{b.name}' cuts through '{a.name}' by {oy:.1f}pt",
                        a.name, b.name, f"Move '{b.name}' above or below '{a.name}'", True))
                elif a.text and b.text:
                    issues.append(Issue(slide_num, Severity.CRITICAL, "overlap",
                        f"'{a.name}' overlaps '{b.name}' by {oy:.1f}pt vertically",
                        a.name, b.name, f"Move lower shape down by {oy + 8:.0f}pt", True))
        if not a.is_background:
            if a.bbox.right > slide_w + 2:
                issues.append(Issue(slide_num, Severity.CRITICAL, "overflow",
                    f"'{a.name}' extends {a.bbox.right - slide_w:.1f}pt beyond right edge",
                    a.name, fix_description="Move left or reduce width", auto_fixable=True))
            if a.bbox.bottom > slide_h + 2:
                issues.append(Issue(slide_num, Severity.CRITICAL, "overflow",
                    f"'{a.name}' extends {a.bbox.bottom - slide_h:.1f}pt beyond bottom edge",
                    a.name, fix_description="Move up or reduce height", auto_fixable=True))
    return issues


# ── Phase 2: Text Truth Check ────────────────────────────────────────────────

def phase2_text_truth_check(shapes, slide_num, slide_w, slide_h):
    issues = []
    content_shapes = [s for s in shapes if not s.is_background]
    for s in content_shapes:
        if s.has_text_overflow:
            overflow = s.estimated_text_height - s.bbox.height
            issues.append(Issue(slide_num, Severity.WARNING, "wrapping",
                f"'{s.name}' text overflows box: stored={s.bbox.height:.0f}pt, "
                f"estimated={s.estimated_text_height:.0f}pt (+{overflow:.0f}pt)",
                s.name, fix_description=f"Resize height to {s.estimated_text_height:.0f}pt",
                auto_fixable=True))
    for i, a in enumerate(content_shapes):
        for b in content_shapes[i + 1:]:
            if "scrim" in a.name.lower() or "overlay" in a.name.lower(): continue
            if "scrim" in b.name.lower() or "overlay" in b.name.lower(): continue
            if not a.has_text_overflow and not b.has_text_overflow: continue
            ox, oy = calc_overlap_with_estimated_height(a, b)
            stored_ox, stored_oy = calc_overlap(a.bbox, b.bbox)
            if ox > 0 and oy > 2 and stored_oy <= 2:
                issues.append(Issue(slide_num, Severity.CRITICAL, "overlap",
                    f"HIDDEN overlap: '{a.name}' true bottom overlaps '{b.name}' "
                    f"by ~{oy:.0f}pt (stored dimensions show no overlap)",
                    a.name, b.name,
                    f"Resize text box, then move '{b.name}' down by {oy + 8:.0f}pt", True))
    return issues


# ── Phase 3: Visual & Readability Audit ──────────────────────────────────────

def get_luminance(hex_color: str) -> float:
    if not hex_color or len(hex_color) < 6: return 0.5
    hex_color = hex_color.replace("#", "").upper()
    try: r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    except ValueError: return 0.5
    def linearize(c):
        c = c / 255.0
        return c / 12.92 if c <= 0.04045 else ((c + 0.055) / 1.055) ** 2.4
    return 0.2126 * linearize(r) + 0.7152 * linearize(g) + 0.0722 * linearize(b)

def contrast_ratio(lum1, lum2):
    return (max(lum1, lum2) + 0.05) / (min(lum1, lum2) + 0.05)

def phase3_visual_audit(shapes, slide_num, has_bg_image):
    issues = []
    for s in shapes:
        if not s.text or s.is_background: continue
        if s.font_size_pt > 0:
            if s.font_size_pt < MIN_FONT_SIZE:
                issues.append(Issue(slide_num, Severity.CRITICAL, "font_size",
                    f"'{s.name}' font is {s.font_size_pt}pt (min {MIN_FONT_SIZE}pt)",
                    s.name, fix_description=f"Increase to {MIN_FONT_SIZE}pt", auto_fixable=True))
            elif s.font_size_pt < MIN_BODY_FONT_SIZE and not s.is_page_furniture:
                issues.append(Issue(slide_num, Severity.WARNING, "font_size",
                    f"'{s.name}' font is {s.font_size_pt}pt (body recommended >={MIN_BODY_FONT_SIZE}pt)",
                    s.name, fix_description=f"Consider increasing to {MIN_BODY_FONT_SIZE}pt"))
        if s.font_color and has_bg_image:
            text_lum = get_luminance(s.font_color)
            bg_lum = 0.05  # assume dark image background
            cr = contrast_ratio(text_lum, bg_lum)
            if cr < 3.0:
                issues.append(Issue(slide_num, Severity.CRITICAL, "contrast",
                    f"'{s.name}' very low contrast ({cr:.1f}:1, #{s.font_color} on dark bg)",
                    s.name, fix_description="Add scrim or change color", auto_fixable=True))
            elif cr < 4.5:
                issues.append(Issue(slide_num, Severity.WARNING, "contrast",
                    f"'{s.name}' low contrast ({cr:.1f}:1, #{s.font_color})", s.name,
                    fix_description="Consider adding scrim overlay"))
        if s.font_color and has_bg_image:
            r_val = int(s.font_color[:2], 16) if len(s.font_color) >= 6 else 0
            g_val = int(s.font_color[2:4], 16) if len(s.font_color) >= 6 else 0
            b_val = int(s.font_color[4:6], 16) if len(s.font_color) >= 6 else 0
            if r_val > 150 and g_val < 80 and b_val < 80:
                issues.append(Issue(slide_num, Severity.WARNING, "contrast",
                    f"'{s.name}' uses red text (#{s.font_color}) on dark bg — low luminance",
                    s.name, fix_description="Consider lighter red or add scrim"))
    return issues


# ── Phase 4: Layout Consistency ──────────────────────────────────────────────

def phase4_consistency(all_slides_shapes, slide_w, slide_h):
    issues = []
    for slide_num, shapes in all_slides_shapes:
        for s in shapes:
            if s.is_background or "scrim" in s.name.lower() or "overlay" in s.name.lower(): continue
            if s.bbox.left < MIN_MARGIN_PT and s.bbox.width < slide_w * 0.8:
                issues.append(Issue(slide_num, Severity.INFO, "margin",
                    f"'{s.name}' is {s.bbox.left:.0f}pt from left edge (min {MIN_MARGIN_PT}pt)", s.name))
            if slide_w - s.bbox.right < MIN_MARGIN_PT and s.bbox.width < slide_w * 0.8:
                issues.append(Issue(slide_num, Severity.INFO, "margin",
                    f"'{s.name}' is {slide_w - s.bbox.right:.0f}pt from right edge", s.name))
    for slide_num, shapes in all_slides_shapes:
        text_shapes = sorted(
            [s for s in shapes if s.text and not s.is_background and not s.is_page_furniture],
            key=lambda s: s.bbox.top)
        if len(text_shapes) < 3: continue
        gaps = []
        for i in range(len(text_shapes) - 1):
            gap = text_shapes[i + 1].bbox.top - text_shapes[i].bbox.bottom
            gaps.append((text_shapes[i].name, text_shapes[i + 1].name, gap))
        if gaps:
            avg_gap = sum(g[2] for g in gaps) / len(gaps)
            for name_a, name_b, gap in gaps:
                if abs(gap - avg_gap) > 15 and len(gaps) > 2:
                    issues.append(Issue(slide_num, Severity.INFO, "spacing",
                        f"Uneven gap between '{name_a}' and '{name_b}': {gap:.0f}pt (avg {avg_gap:.0f}pt)"))
    for slide_num, shapes in all_slides_shapes:
        for s in shapes:
            if s.shape_type == "textbox" and not s.text.strip() and not s.is_decorative and not s.has_fill:
                issues.append(Issue(slide_num, Severity.INFO, "unused",
                    f"'{s.name}' is empty text shape — consider deleting",
                    s.name, fix_description="Delete empty placeholder", auto_fixable=True))
    return issues


# ── Phase 5: Composition Coverage ─────────────────────────────────────────────

def phase5_composition_coverage(shapes, slide_num, slide_w, slide_h):
    """
    Check whether overlay shapes on slides with full-bleed BG images cover
    too much of the slide, blocking the image's visual content.
    """
    issues = []
    slide_area = slide_w * slide_h
    if slide_area <= 0:
        return issues

    # Find full-bleed background image (picture at ~0,0 covering >=90% of slide)
    bg_image = None
    for s in shapes:
        if s.shape_type == "image" and s.is_background:
            if s.bbox.left <= 5 and s.bbox.top <= 5:  # near origin
                bg_image = s
                break

    if bg_image is None:
        return issues

    # Find overlay shapes: non-picture, non-background shapes with a solid fill
    # Exclude text-only shapes (textboxes without fill) and the BG image itself
    overlay_shapes = []
    for s in shapes:
        if s is bg_image or s.is_background:
            continue
        if s.shape_type == "image":
            continue
        if not s.has_fill:
            continue
        # Include filled shapes (rectangles, rounded rects, cards, scrims, etc.)
        overlay_shapes.append(s)

    if not overlay_shapes:
        return issues

    # Check individual overlay coverage
    total_overlay_area = 0
    for s in overlay_shapes:
        shape_area = s.bbox.width * s.bbox.height
        pct = (shape_area / slide_area) * 100
        total_overlay_area += shape_area

        if pct > 30:
            issues.append(Issue(
                slide_num, Severity.WARNING, "composition",
                f"'{s.name}' covers {pct:.0f}% of slide with full-bleed BG image "
                f"— may block key visual content. Consider reducing size or removing.",
                s.name))

    # Check combined overlay coverage
    total_pct = (total_overlay_area / slide_area) * 100
    if total_pct > 50:
        issues.append(Issue(
            slide_num, Severity.CRITICAL, "composition",
            f"Combined overlay shapes cover {total_pct:.0f}% of slide "
            f"— BG image is mostly hidden. Consider removing overlays."))

    return issues


# ── Fix Engine ────────────────────────────────────────────────────────────────

REFLOW_GAP_PT = 4  # Minimum gap between shapes after reflow nudge

def _is_bg_image(shape, slide_w_pt, slide_h_pt):
    """Check if a shape is a full-slide background image."""
    return (hasattr(shape, 'image') and shape.width and shape.height
            and emu_to_pt(shape.width) >= slide_w_pt * 0.9
            and emu_to_pt(shape.height) >= slide_h_pt * 0.9)

def fix_font_size(shape_info, min_size_pt):
    if not shape_info._shape_ref.has_text_frame: return
    from lxml import etree
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    min_sz = int(min_size_pt * 100)  # hundredths of a point
    for r_el in shape_info._shape_ref._element.findall('.//a:r', nsmap):
        rPr = r_el.find('a:rPr', nsmap)
        if rPr is not None:
            sz = rPr.get('sz')
            if sz and int(sz) < min_sz:
                rPr.set('sz', str(min_sz))

def add_scrim_overlay(slide, left_pt, top_pt, width_pt, height_pt, alpha_pct=45):
    from pptx.oxml.ns import qn
    spTree = slide.shapes._spTree
    sp_xml = (
        f'<p:sp xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<p:nvSpPr><p:cNvPr id="0" name="Scrim Overlay"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>'
        f'<p:spPr>'
        f'<a:xfrm><a:off x="{int(Pt(left_pt))}" y="{int(Pt(top_pt))}"/>'
        f'<a:ext cx="{int(Pt(width_pt))}" cy="{int(Pt(height_pt))}"/></a:xfrm>'
        f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
        f'<a:solidFill><a:srgbClr val="000000">'
        f'<a:alpha val="{alpha_pct * 1000}"/>'
        f'</a:srgbClr></a:solidFill>'
        f'<a:ln><a:noFill/></a:ln>'
        f'</p:spPr>'
        f'<p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr/></a:p></p:txBody>'
        f'</p:sp>')
    sp_element = etree.fromstring(sp_xml)
    pic_elements = spTree.findall(qn("p:pic"))
    if pic_elements: pic_elements[0].addnext(sp_element)
    else:
        children = list(spTree)
        if len(children) >= 2: children[1].addnext(sp_element)
        else: spTree.append(sp_element)

def _has_shape_fill(shape):
    """Check if shape has a visible fill via raw XML (no python-pptx side effects)."""
    try:
        nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
        sp_el = shape._element
        spPr = sp_el.find('.//a:spPr', nsmap)
        if spPr is not None:
            if (spPr.find('a:solidFill', nsmap) is not None
                or spPr.find('a:gradFill', nsmap) is not None
                or spPr.find('a:pattFill', nsmap) is not None
                or spPr.find('a:blipFill', nsmap) is not None):
                return True
        # Check style-referenced fill (fillRef idx > 0)
        fillRef = sp_el.find('.//a:fillRef', nsmap)
        if fillRef is not None and fillRef.get('idx', '0') != '0':
            return True
    except: pass
    return False

def _find_background_partners(shape, all_shapes, slide_w_pt, slide_h_pt):
    """
    Find filled background shapes that visually sit behind a text shape.
    A background partner is a shape that:
    - Has a visible fill (solid, gradient, pattern, or theme-referenced)
    - Has NO text (or only whitespace)
    - Spatially contains or mostly overlaps the text shape
    Returns list of matching shapes (usually 0 or 1).
    """
    if shape.top is None or shape.left is None:
        return []
    s_left = shape.left or 0
    s_top = shape.top or 0
    s_right = s_left + (shape.width or 0)
    s_bottom = s_top + (shape.height or 0)
    partners = []
    for other in all_shapes:
        if other is shape:
            continue
        if other.top is None or other.left is None:
            continue
        # Must be a filled shape with no text
        if not _has_shape_fill(other):
            continue
        if other.has_text_frame and other.text_frame.text.strip():
            continue
        # Skip full-slide background images
        if _is_bg_image(other, slide_w_pt, slide_h_pt):
            continue
        o_left = other.left or 0
        o_top = other.top or 0
        o_right = o_left + (other.width or 0)
        o_bottom = o_top + (other.height or 0)
        # The background shape should be a card-sized element that contains the text shape,
        # NOT a large page-spanning panel. Filter by:
        # 1. Text shape's center must be inside background shape
        # 2. Background shape must be card-sized (not more than 3x the text shape's dimensions)
        s_cx = (s_left + s_right) / 2
        s_cy = (s_top + s_bottom) / 2
        if o_left <= s_cx <= o_right and o_top <= s_cy <= o_bottom:
            s_h = (shape.height or 0)
            s_w = (shape.width or 0)
            o_h = (other.height or 0)
            o_w = (other.width or 0)
            # Skip if background is vastly larger than the text shape (page-level panel)
            if s_h > 0 and o_h > s_h * 4:
                continue
            if s_w > 0 and o_w > s_w * 4:
                continue
            partners.append(other)
    return partners


def _reflow_slide(slide, slide_w_pt, slide_h_pt):
    """
    Iterative vertical reflow: after shapes have been resized, resolve all
    overlaps by nudging lower shapes down. Only moves text-bearing shapes —
    filled background shapes (cards, panels) are anchored in place.
    Returns number of shapes moved.
    """
    movable = []
    for shape in slide.shapes:
        if _is_bg_image(shape, slide_w_pt, slide_h_pt):
            continue
        if shape.top is None:
            continue
        # Skip filled shapes (card backgrounds, decorative panels) — only reflow text content
        if _has_shape_fill(shape) and not (shape.has_text_frame and shape.text_frame.text.strip()):
            continue
        movable.append(shape)

    move_count = 0
    changed = True
    iterations = 0
    while changed and iterations < 20:
        changed = False
        iterations += 1
        movable.sort(key=lambda s: (s.top, s.left))
        for i, upper in enumerate(movable):
            upper_bottom = upper.top + upper.height
            for j in range(i + 1, len(movable)):
                lower = movable[j]
                if lower.top is None or upper.top is None:
                    continue
                # Check horizontal overlap
                u_left = upper.left or 0
                u_right = u_left + (upper.width or 0)
                l_left = lower.left or 0
                l_right = l_left + (lower.width or 0)
                h_overlap = min(u_right, l_right) - max(u_left, l_left)
                if h_overlap <= 0:
                    continue
                # If they vertically overlap, push lower shape down
                gap_emu = lower.top - upper_bottom
                if gap_emu < 0:
                    nudge = -gap_emu + Pt(REFLOW_GAP_PT)
                    lower.top = lower.top + nudge
                    # Move background partners with the text shape
                    partners = _find_background_partners(lower, slide.shapes, slide_w_pt, slide_h_pt)
                    for partner in partners:
                        partner.top = partner.top + nudge
                    changed = True
                    move_count += 1
    return move_count

def auto_fix_issues(prs, report, all_slides_shapes):
    """
    Apply automatic fixes with proper vertical reflow.

    Strategy:
    1. Fix font sizes (safe, no layout side effects)
    2. Delete empty placeholders (safe)
    3. Resize overflowing text boxes (only significant overflows > MIN_OVERFLOW_FIX_PT)
    4. Run iterative vertical reflow per slide to resolve any new overlaps
       caused by the resizing — this is the critical step the old engine missed.
    """
    fix_count = 0
    slide_w_pt = emu_to_pt(prs.slide_width)
    slide_h_pt = emu_to_pt(prs.slide_height)

    shape_lookup = {}
    for slide_num, shapes in all_slides_shapes:
        for s in shapes:
            shape_lookup[(slide_num, s.name)] = s

    # Pass 1: Font size fixes (no layout impact)
    for issue in report.issues:
        if not issue.auto_fixable: continue
        if issue.category == "font_size" and issue.severity == Severity.CRITICAL:
            key = (issue.slide_num, issue.shape_a)
            if key in shape_lookup:
                fix_font_size(shape_lookup[key], MIN_FONT_SIZE)
                fix_count += 1

    # Pass 2: Delete empty placeholders (no layout impact)
    for issue in report.issues:
        if not issue.auto_fixable: continue
        if issue.category == "unused":
            key = (issue.slide_num, issue.shape_a)
            if key in shape_lookup:
                sp = shape_lookup[key]._shape_ref._element
                sp.getparent().remove(sp)
                fix_count += 1

    # Pass 3: Fix overflowing text boxes — WIDEN first, then HEIGHT as fallback
    # Many overflows are caused by unwanted text wrapping (box too narrow),
    # not insufficient height. Widening avoids cascading reflow entirely.
    slides_needing_reflow = set()
    for issue in report.issues:
        if not issue.auto_fixable: continue
        if issue.category == "wrapping":
            key = (issue.slide_num, issue.shape_a)
            if key not in shape_lookup:
                continue
            si = shape_lookup[key]
            growth = si.estimated_text_height - si.bbox.height
            if growth <= MIN_OVERFLOW_FIX_PT:
                continue  # Skip small overflows — PowerPoint handles them fine

            # Strategy A: Try widening the box to reduce line wraps
            # Calculate width needed for text to fit in fewer lines
            widened = False
            if si.text and si.font_size_pt > 0:
                ratio = CHAR_WIDTH_RATIO.get(si.font_name, CHAR_WIDTH_RATIO["default"])
                if si.font_bold:
                    ratio *= 1.08
                avg_char_w = si.font_size_pt * ratio
                line_h = si.font_size_pt * LINE_HEIGHT_RATIO

                # How many lines would fit in the current box height?
                max_lines_in_box = max(1, int(si.bbox.height / line_h))

                # What width would we need so text fits in max_lines_in_box lines?
                paragraphs = si.text.split("\n")
                longest_para_chars = max((len(p) for p in paragraphs), default=0)
                needed_width_pt = (longest_para_chars * avg_char_w / max_lines_in_box) + 20  # +20 for padding

                # Available space: can we widen rightward without going past slide edge - margin?
                max_right_pt = slide_w_pt - MIN_MARGIN_PT
                current_right_pt = si.bbox.left + si.bbox.width
                available_width_pt = max_right_pt - si.bbox.left

                if needed_width_pt <= available_width_pt and needed_width_pt > si.bbox.width:
                    # Check that the new width actually eliminates the overflow
                    new_est_height = estimate_text_height(
                        si.text, si.font_size_pt, si.font_name,
                        needed_width_pt - 14, si.font_bold)
                    if new_est_height <= si.bbox.height * 1.05:
                        si._shape_ref.width = Pt(needed_width_pt)
                        widened = True
                        fix_count += 1
                        print(f"  Widened '{si.name}' on slide {issue.slide_num}: "
                              f"{si.bbox.width:.0f} -> {needed_width_pt:.0f}pt (avoids wrapping)")

            # Strategy B: Fall back to height increase if widening didn't work
            if not widened:
                si._shape_ref.height = Pt(si.estimated_text_height)
                fix_count += 1
                slides_needing_reflow.add(issue.slide_num)

    # Pass 4: Fix structural overlaps by nudging the lower shape down
    for issue in report.issues:
        if not issue.auto_fixable: continue
        if issue.category != "overlap": continue
        if not issue.shape_a or not issue.shape_b: continue
        key_a = (issue.slide_num, issue.shape_a)
        key_b = (issue.slide_num, issue.shape_b)
        if key_a not in shape_lookup or key_b not in shape_lookup: continue
        a = shape_lookup[key_a]
        b = shape_lookup[key_b]
        # Skip decorative line overlaps (thin dividers) — these are intentional design elements
        if a.is_decorative or b.is_decorative: continue
        # Only fix text-on-text overlaps
        if not (a.text and b.text): continue
        # Determine which is lower (the one to nudge)
        upper, lower = (a, b) if a.bbox.top <= b.bbox.top else (b, a)
        upper_bottom = upper.bbox.top + max(upper.bbox.height, upper.estimated_text_height)
        overlap_amount = upper_bottom - lower.bbox.top
        if overlap_amount <= 0: continue
        # Skip minor overlaps (< MIN_OVERFLOW_FIX_PT) — PowerPoint handles these,
        # and fixing them can trigger cascade reflow that destroys card layouts
        if overlap_amount < MIN_OVERFLOW_FIX_PT: continue
        nudge_pt = overlap_amount + REFLOW_GAP_PT
        nudge_emu = Pt(nudge_pt)
        lower._shape_ref.top = lower._shape_ref.top + nudge_emu
        fix_count += 1
        slides_needing_reflow.add(issue.slide_num)
        print(f"  Nudged '{lower.name}' down {nudge_pt:.0f}pt on slide {issue.slide_num} to fix overlap")
        # Move background partner(s) with the text shape so cards stay intact
        slide = prs.slides[issue.slide_num - 1]
        partners = _find_background_partners(lower._shape_ref, slide.shapes, slide_w_pt, slide_h_pt)
        for partner in partners:
            partner.top = partner.top + nudge_emu
            print(f"    Also moved background '{partner.name}' to keep card intact")

    # Pass 5: Vertical reflow on affected slides
    for slide_num in sorted(slides_needing_reflow):
        slide = prs.slides[slide_num - 1]
        moved = _reflow_slide(slide, slide_w_pt, slide_h_pt)
        if moved:
            print(f"  Slide {slide_num}: reflowed {moved} shape(s) to resolve overlaps")

    return fix_count


# ── Main Auditor Class ───────────────────────────────────────────────────────

class PptxAuditor:
    """
    Main entry point. Usage:
        auditor = PptxAuditor("my_deck.pptx")
        report = auditor.run_full_audit()
        print(report)
        auditor.fix_all(report)
        auditor.save("my_deck_fixed.pptx")
    """
    def __init__(self, filepath):
        self.filepath = filepath
        self.prs = Presentation(filepath)
        self.slide_width_pt = emu_to_pt(self.prs.slide_width)
        self.slide_height_pt = emu_to_pt(self.prs.slide_height)
        self.all_slides_shapes = []

    def _extract_all_shapes(self):
        self.all_slides_shapes = []
        for idx, slide in enumerate(self.prs.slides):
            slide_num = idx + 1
            shapes = []
            for shape in slide.shapes:
                try: shapes.append(extract_shape_info(shape, self.slide_width_pt, self.slide_height_pt))
                except Exception as e: print(f"  Warning: shape '{shape.name}' slide {slide_num}: {e}")
            self.all_slides_shapes.append((slide_num, shapes))

    def run_full_audit(self):
        print(f"Auditing: {self.filepath}")
        print(f"Slide size: {self.slide_width_pt:.0f} x {self.slide_height_pt:.0f} pt")
        self._extract_all_shapes()
        report = AuditReport(os.path.basename(self.filepath), len(self.all_slides_shapes),
                             self.slide_width_pt, self.slide_height_pt)
        for slide_num, shapes in self.all_slides_shapes:
            has_bg_image = any(s.is_background for s in shapes)
            report.issues.extend(phase1_structural_scan(shapes, slide_num,
                                 self.slide_width_pt, self.slide_height_pt))
            report.issues.extend(phase2_text_truth_check(shapes, slide_num,
                                 self.slide_width_pt, self.slide_height_pt))
            report.issues.extend(phase3_visual_audit(shapes, slide_num, has_bg_image))
            report.issues.extend(phase5_composition_coverage(shapes, slide_num,
                                 self.slide_width_pt, self.slide_height_pt))
        report.issues.extend(phase4_consistency(self.all_slides_shapes,
                             self.slide_width_pt, self.slide_height_pt))
        print(f"\nAudit complete: {report.critical_count} critical, "
              f"{report.warning_count} warnings, {report.info_count} info")
        return report

    def fix_all(self, report):
        count = auto_fix_issues(self.prs, report, self.all_slides_shapes)
        print(f"Applied {count} automatic fixes (with vertical reflow)")
        return count

    def fix_and_verify(self, max_passes=3):
        """
        Fix-verify loop: audit -> fix -> re-audit -> repeat until stable.
        Returns (final_report, total_fixes, passes_used).

        Each pass:
        1. Re-extracts shape info (positions may have changed from prior fixes)
        2. Runs full audit
        3. Applies auto-fixes
        4. If no fixes were applied or no auto-fixable criticals remain, stop

        Max passes prevents infinite loops if fixes keep creating new issues.
        """
        total_fixes = 0
        for pass_num in range(1, max_passes + 1):
            print(f"\n{'='*50}")
            print(f"Pass {pass_num}/{max_passes}")
            print(f"{'='*50}")

            # Re-extract shapes (positions changed from prior fixes)
            self._extract_all_shapes()

            # Run audit
            report = AuditReport(os.path.basename(self.filepath), len(self.all_slides_shapes),
                                 self.slide_width_pt, self.slide_height_pt)
            for slide_num, shapes in self.all_slides_shapes:
                has_bg_image = any(s.is_background for s in shapes)
                report.issues.extend(phase1_structural_scan(shapes, slide_num,
                                     self.slide_width_pt, self.slide_height_pt))
                report.issues.extend(phase2_text_truth_check(shapes, slide_num,
                                     self.slide_width_pt, self.slide_height_pt))
                report.issues.extend(phase3_visual_audit(shapes, slide_num, has_bg_image))
                report.issues.extend(phase5_composition_coverage(shapes, slide_num,
                                     self.slide_width_pt, self.slide_height_pt))
            report.issues.extend(phase4_consistency(self.all_slides_shapes,
                                 self.slide_width_pt, self.slide_height_pt))

            auto_fixable = [i for i in report.issues if i.auto_fixable]
            print(f"Audit: {report.critical_count} critical, {report.warning_count} warnings, "
                  f"{report.info_count} info ({len(auto_fixable)} auto-fixable)")

            if not auto_fixable:
                print("No auto-fixable issues remain. Done.")
                break

            # Apply fixes
            fixes = auto_fix_issues(self.prs, report, self.all_slides_shapes)
            total_fixes += fixes
            print(f"Applied {fixes} fixes in pass {pass_num}")

            if fixes == 0:
                print("No fixes applied (all remaining issues need manual review). Done.")
                break

        # Final verification audit
        print(f"\n{'='*50}")
        print(f"Final verification (after {pass_num} pass{'es' if pass_num > 1 else ''})")
        print(f"{'='*50}")
        self._extract_all_shapes()
        final_report = AuditReport(os.path.basename(self.filepath), len(self.all_slides_shapes),
                                   self.slide_width_pt, self.slide_height_pt)
        for slide_num, shapes in self.all_slides_shapes:
            has_bg_image = any(s.is_background for s in shapes)
            final_report.issues.extend(phase1_structural_scan(shapes, slide_num,
                                       self.slide_width_pt, self.slide_height_pt))
            final_report.issues.extend(phase2_text_truth_check(shapes, slide_num,
                                       self.slide_width_pt, self.slide_height_pt))
            final_report.issues.extend(phase3_visual_audit(shapes, slide_num, has_bg_image))
            final_report.issues.extend(phase5_composition_coverage(shapes, slide_num,
                                       self.slide_width_pt, self.slide_height_pt))
        final_report.issues.extend(phase4_consistency(self.all_slides_shapes,
                                   self.slide_width_pt, self.slide_height_pt))

        remaining_auto = [i for i in final_report.issues if i.auto_fixable]
        print(f"Final: {final_report.critical_count} critical, {final_report.warning_count} warnings, "
              f"{final_report.info_count} info")
        print(f"Total fixes applied: {total_fixes} across {pass_num} pass(es)")
        if remaining_auto:
            print(f"Note: {len(remaining_auto)} auto-fixable issues remain (below threshold or marginal)")
        else:
            print("All auto-fixable issues resolved.")

        return final_report, total_fixes, pass_num

    def save(self, output_path):
        self.prs.save(output_path)
        print(f"Saved to: {output_path}")


if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("Usage: python pptx_audit.py <file.pptx> [--fix] [--output fixed.pptx]"); sys.exit(1)
    filepath = sys.argv[1]
    do_fix = "--fix" in sys.argv
    output = None
    if "--output" in sys.argv:
        idx = sys.argv.index("--output")
        if idx + 1 < len(sys.argv): output = sys.argv[idx + 1]
    auditor = PptxAuditor(filepath)
    if do_fix:
        final_report, total_fixes, passes = auditor.fix_and_verify()
        print("\n" + str(final_report))
        auditor.save(output or filepath.replace(".pptx", "_fixed.pptx"))
    else:
        report = auditor.run_full_audit()
        print("\n" + str(report))
